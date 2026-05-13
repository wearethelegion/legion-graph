"""
Document text extraction layer for multiple formats.

Supports PDF, DOCX, Markdown, TXT, and code files with structure preservation.
"""

import os
from pathlib import Path
from typing import Optional

import pypdf
from docx import Document
from pptx import Presentation


def detect_file_type(filename: str) -> str:
    """
    Detect file type from extension.

    Args:
        filename: Path or filename to analyze

    Returns:
        One of: "pdf", "docx", "markdown", "text", "code"

    Raises:
        ValueError: If file extension is not supported
    """
    ext = Path(filename).suffix.lower()

    # PDF files
    if ext == ".pdf":
        return "pdf"

    # DOCX files
    if ext in [".docx", ".doc"]:
        return "docx"

    # PowerPoint files
    if ext in [".pptx", ".ppt"]:
        return "pptx"

    # Markdown files
    if ext in [".md", ".markdown"]:
        return "markdown"

    # Plain text files
    if ext in [".txt", ".text"]:
        return "text"

    # Code files
    code_extensions = {
        ".py", ".js", ".ts", ".java", ".go", ".rs",
        ".rb", ".php", ".c", ".cpp", ".h", ".hpp",
        ".cs", ".swift", ".kt", ".scala", ".r", ".m"
    }
    if ext in code_extensions:
        return "code"

    raise ValueError(f"Unsupported file extension: {ext}")


def detect_code_language(filename: str) -> str:
    """
    Detect programming language from file extension.

    Args:
        filename: Path or filename to analyze

    Returns:
        Language name (e.g., "python", "javascript", "typescript")
    """
    ext = Path(filename).suffix.lower()

    language_map = {
        ".py": "python",
        ".js": "javascript",
        ".ts": "typescript",
        ".java": "java",
        ".go": "go",
        ".rs": "rust",
        ".rb": "ruby",
        ".php": "php",
        ".c": "c",
        ".cpp": "cpp",
        ".h": "c_header",
        ".hpp": "cpp_header",
        ".cs": "csharp",
        ".swift": "swift",
        ".kt": "kotlin",
        ".scala": "scala",
        ".r": "r",
        ".m": "objective_c",
    }

    return language_map.get(ext, "unknown")


async def extract_text(file_path: str) -> str:
    """
    Main dispatcher function for text extraction.

    Routes to appropriate extractor based on file type.

    Args:
        file_path: Path to the file to extract text from

    Returns:
        Extracted text content

    Raises:
        ValueError: If file type is not supported
        FileNotFoundError: If file does not exist
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    file_type = detect_file_type(file_path)

    if file_type == "pdf":
        return await extract_text_from_pdf(file_path)
    elif file_type == "docx":
        return await extract_text_from_docx(file_path)
    elif file_type == "pptx":
        return await extract_text_from_pptx(file_path)
    elif file_type in ["markdown", "text"]:
        return await extract_text_from_text(file_path)
    elif file_type == "code":
        return await extract_text_from_code(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


async def extract_text_from_pdf(file_path: str) -> str:
    """
    Extract text from PDF preserving structure.

    Uses pypdf library to extract text from all pages.
    Adds page markers for reference: [PAGE N]

    Args:
        file_path: Path to PDF file

    Returns:
        Extracted text with page markers

    Raises:
        Exception: If PDF reading fails
    """
    text_parts = []

    try:
        with open(file_path, 'rb') as f:
            pdf = pypdf.PdfReader(f)

            for page_num, page in enumerate(pdf.pages, start=1):
                page_text = page.extract_text()

                # Add page markers for later reference
                text_parts.append(f"\n[PAGE {page_num}]\n{page_text}")

        return "\n".join(text_parts)

    except Exception as e:
        raise Exception(f"Failed to extract text from PDF {file_path}: {str(e)}")


async def extract_text_from_docx(file_path: str) -> str:
    """
    Extract text from DOCX preserving structure.

    Uses python-docx library to extract paragraphs.
    Converts heading styles to markdown format.

    Args:
        file_path: Path to DOCX file

    Returns:
        Extracted text with markdown-style headings

    Raises:
        Exception: If DOCX reading fails
    """
    try:
        doc = Document(file_path)
        text_parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                # Preserve heading levels by converting to markdown
                if para.style.name.startswith('Heading'):
                    # Extract heading level (e.g., 'Heading 1' -> 1)
                    try:
                        level = int(para.style.name.split()[-1])
                        text_parts.append(f"\n{'#' * level} {para.text}\n")
                    except (ValueError, IndexError):
                        # If we can't parse the level, treat as regular text
                        text_parts.append(para.text)
                else:
                    text_parts.append(para.text)

        return "\n".join(text_parts)

    except Exception as e:
        raise Exception(f"Failed to extract text from DOCX {file_path}: {str(e)}")


async def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract text from PowerPoint (PPTX/PPT) files.

    Uses python-pptx library to extract text from all slides.
    Includes slide numbers and preserves text from shapes, tables, and notes.

    Args:
        file_path: Path to PowerPoint file

    Returns:
        Extracted text with slide markers

    Raises:
        Exception: If PowerPoint reading fails
    """
    try:
        prs = Presentation(file_path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []

            # Add slide marker
            slide_text.append(f"\n[SLIDE {slide_num}]\n")

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                        if row_text:
                            slide_text.append(row_text)

            # Extract notes if present
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_text.append(f"\n[NOTES]\n{notes_text}")

            text_parts.append("\n".join(slide_text))

        return "\n".join(text_parts)

    except Exception as e:
        raise Exception(f"Failed to extract text from PowerPoint {file_path}: {str(e)}")


async def extract_text_from_text(file_path: str) -> str:
    """
    Extract text from plain text or markdown files.

    Simple UTF-8 text file read.

    Args:
        file_path: Path to text file

    Returns:
        File content as-is

    Raises:
        Exception: If file reading fails
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    except UnicodeDecodeError:
        # Try with different encoding if UTF-8 fails
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            raise Exception(f"Failed to read text file {file_path}: {str(e)}")

    except Exception as e:
        raise Exception(f"Failed to extract text from file {file_path}: {str(e)}")


async def extract_text_from_code(file_path: str) -> str:
    """
    Extract text from code files with structure preservation.

    Reads code file and adds metadata header with filename and language.

    Args:
        file_path: Path to code file

    Returns:
        Code content with metadata header

    Raises:
        Exception: If code file reading fails
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            code = f.read()

        # Detect language
        language = detect_code_language(file_path)
        filename = os.path.basename(file_path)

        # Preserve code structure with markers
        return f"""[CODE FILE: {filename}]
[LANGUAGE: {language}]

{code}"""

    except UnicodeDecodeError:
        # Try with different encoding if UTF-8 fails
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                code = f.read()

            language = detect_code_language(file_path)
            filename = os.path.basename(file_path)

            return f"""[CODE FILE: {filename}]
[LANGUAGE: {language}]

{code}"""
        except Exception as e:
            raise Exception(f"Failed to read code file {file_path}: {str(e)}")

    except Exception as e:
        raise Exception(f"Failed to extract text from code file {file_path}: {str(e)}")
